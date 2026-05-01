"""Stakeholder-ready DOCX + PPTX export.

Turns the same corpus the existing markdown export reads (`topic_insights`
+ `posts` + `topic_posts` + `graph_nodes`) into formats people actually
present from:

  - DOCX (python-docx) — the long-form research brief: cover page, exec
    summary, painpoints with cited evidence quotes, competitor teardown,
    feature roadmap, citation appendix. This is the doc you email an
    investor or a CTO.
  - PPTX (python-pptx) — a 12-15 slide pitch deck: TL;DR, problem, top
    painpoints with the strongest quote per slide, competitor matrix,
    feature plan, ask. This is the doc you screen-share in a meeting.

Why python-docx + python-pptx and nothing else:
  - Both are pure-Python, zero system deps (no LibreOffice, no pandoc),
    work fine inside a PyInstaller sidecar.
  - Both are battle-tested (10+ years), MIT-licensed, ~5 MB each.
  - Output is real .docx / .pptx (Office Open XML) — opens in Word,
    Google Docs, Pages, Keynote, PowerPoint, LibreOffice without any
    rendering surprises.
  - PDF: deliberately not bundled. python-docx2pdf needs Word/LibreOffice
    on the host; weasyprint needs cairo/pango. Users who need PDF can
    open the .docx in Word and Save-as-PDF in 2 clicks. Adding 80 MB of
    rendering libraries to the sidecar is not worth it.

The module degrades gracefully: if the optional packages aren't
installed, every public function returns a structured error dict instead
of raising — so the MCP tool can tell the user exactly what to install.
"""
from __future__ import annotations

import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from ..core.db import get_db


# ─── Soft-import optional deps ────────────────────────────────────────────
# We don't want to crash the MCP server at import time on a clean install
# where docx/pptx aren't installed yet. The functions check `_DOCX_OK` /
# `_PPTX_OK` at call time and return a clear install hint.

try:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor
    _DOCX_OK = True
except Exception as _e:
    Document = None  # type: ignore
    WD_ALIGN_PARAGRAPH = None  # type: ignore
    Inches = Pt = RGBColor = None  # type: ignore
    _DOCX_OK = False
    _DOCX_ERR = str(_e)

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.shapes import MSO_SHAPE
    _PPTX_OK = True
except Exception as _e:
    Presentation = None  # type: ignore
    PptxInches = PptxPt = PptxRGBColor = MSO_SHAPE = None  # type: ignore
    _PPTX_OK = False
    _PPTX_ERR = str(_e)

# Optional: pandoc-backed markdown → DOCX. `pypandoc-binary` ships its own
# pandoc binary (~25 MB) so end users don't need a system pandoc install.
# Plain `pypandoc` works too but requires pandoc on PATH. We try both.
try:
    import pypandoc  # type: ignore
    _PANDOC_OK = True
except Exception as _e:
    pypandoc = None  # type: ignore
    _PANDOC_OK = False
    _PANDOC_ERR = str(_e)


# ─── Data layer — pull a single struct that both renderers consume ───────


def _gather_deck_data(
    topic: str,
    extra_topics: list[str] | None = None,
    max_painpoints: int = 12,
    max_posts_per_painpoint: int = 4,
    max_competitors: int = 10,
) -> dict[str, Any]:
    """Build the data struct that both DOCX + PPTX renderers read.

    Pulls from the SAME tables the markdown exporter reads, so we never
    drift between formats:

      - `topic_insights` — exec summary, governing thought, key findings,
        competitor list (LLM-synthesized; may be missing if user hasn't
        run `research insights` yet — we degrade gracefully).
      - `topic_posts` JOIN `posts` — top engagement quotes for citation.
      - `graph_nodes` (kind='painpoint') — frequency-weighted painpoint
        list with evidence quotes.
      - `posts` filtered to App Store / Play Store — competitor signal
        with star-rating distribution.

    `extra_topics` lets the user widen the corpus (e.g. include the 3
    sibling lending topics for the marketplace deck) without re-running
    a collect.
    """
    db = get_db()

    topics = [topic] + list(extra_topics or [])
    placeholders = ",".join(["?"] * len(topics))

    out: dict[str, Any] = {
        "topic": topic,
        "extra_topics": list(extra_topics or []),
        "generated_at": datetime.now(timezone.utc).isoformat(),
    }

    # 1) Synthesized insights (optional)
    insights: dict[str, Any] = {}
    if "topic_insights" in db.table_names():
        rows = list(db.query(
            "SELECT report_json FROM topic_insights WHERE topic = ?",
            [topic],
        ))
        if rows:
            try:
                insights = json.loads(rows[0]["report_json"] or "{}")
            except Exception:
                insights = {}
    out["insights"] = insights

    # 2) Corpus stats
    out["total_posts"] = list(db.query(
        f"SELECT COUNT(DISTINCT post_id) AS n FROM topic_posts WHERE topic IN ({placeholders})",
        topics,
    ))[0]["n"]

    out["sources"] = list(db.query(
        f"""SELECT p.source_type, COUNT(DISTINCT p.id) AS n
            FROM topic_posts tp JOIN posts p ON p.id = tp.post_id
            WHERE tp.topic IN ({placeholders})
            GROUP BY p.source_type ORDER BY n DESC""",
        topics,
    ))

    out["sub_breakdown"] = list(db.query(
        f"""SELECT p.sub, COUNT(DISTINCT p.id) AS n
            FROM topic_posts tp JOIN posts p ON p.id = tp.post_id
            WHERE tp.topic IN ({placeholders})
            GROUP BY p.sub ORDER BY n DESC LIMIT 20""",
        topics,
    ))

    # 3) Painpoints from graph_nodes
    painpoints: list[dict[str, Any]] = []
    if "graph_nodes" in db.table_names():
        rows = list(db.query(
            f"""SELECT id, label, metadata_json
                FROM graph_nodes
                WHERE topic IN ({placeholders}) AND kind='painpoint'""",
            topics,
        ))
        for r in rows:
            try:
                meta = json.loads(r["metadata_json"] or "{}")
            except Exception:
                meta = {}
            painpoints.append({
                "id": r["id"],
                "label": r["label"],
                "evidence": meta.get("evidence", ""),
                "frequency": int(meta.get("frequency", 0) or 0),
                "importance": float(meta.get("importance", 0) or 0),
                "satisfaction": float(meta.get("satisfaction", 0) or 0),
                "opportunity_score": float(meta.get("opportunity_score", 0) or 0),
                "pain_weight": float(meta.get("pain_weight", 0) or 0),
            })
        painpoints.sort(
            key=lambda p: (p["opportunity_score"], p["frequency"]),
            reverse=True,
        )
    out["painpoints"] = painpoints[:max_painpoints]

    # 4) Per-painpoint citation pulls — take the top engagement quotes
    # whose title or body mentions the painpoint label words.
    for pp in out["painpoints"]:
        terms = [w for w in pp["label"].split() if len(w) > 3][:4]
        if not terms:
            pp["citations"] = []
            continue
        like_clauses = " AND ".join(
            ["(lower(p.title) || ' ' || lower(coalesce(p.selftext,''))) LIKE ?"] * len(terms)
        )
        params = topics + [f"%{w.lower()}%" for w in terms]
        try:
            cited = list(db.query(
                f"""SELECT DISTINCT p.id, p.sub, p.score, p.title, p.permalink,
                                    p.source_type, substr(coalesce(p.selftext,''),1,400) AS body
                    FROM topic_posts tp JOIN posts p ON p.id = tp.post_id
                    WHERE tp.topic IN ({placeholders}) AND {like_clauses}
                    ORDER BY p.score DESC LIMIT {max_posts_per_painpoint}""",
                params,
            ))
        except Exception:
            cited = []
        pp["citations"] = cited

    # 5) Competitor teardown — reviews per app, plus star distribution.
    out["competitors"] = list(db.query(
        f"""SELECT p.sub,
                   COUNT(*) AS reviews,
                   AVG(p.score) AS avg_stars,
                   SUM(CASE WHEN p.score=1 THEN 1 ELSE 0 END) AS one_star,
                   SUM(CASE WHEN p.score=5 THEN 1 ELSE 0 END) AS five_star
            FROM topic_posts tp JOIN posts p ON p.id = tp.post_id
            WHERE tp.topic IN ({placeholders})
              AND p.source_type IN ('appstore','playstore')
            GROUP BY p.sub
            HAVING reviews >= 5
            ORDER BY reviews DESC LIMIT {max_competitors}""",
        topics,
    ))

    # 6) Top engagement quotes (the deck's "voice of customer" slide)
    out["top_quotes"] = list(db.query(
        f"""SELECT DISTINCT p.id, p.sub, p.score, p.num_comments, p.title,
                            p.permalink, substr(coalesce(p.selftext,''),1,300) AS body
            FROM topic_posts tp JOIN posts p ON p.id = tp.post_id
            WHERE tp.topic IN ({placeholders}) AND p.source_type='reddit' AND p.score >= 100
            ORDER BY p.score DESC LIMIT 8""",
        topics,
    ))

    return out


# ─── DOCX renderer ────────────────────────────────────────────────────────


def _severity_for_score(opp: float) -> str:
    if opp >= 11:
        return "high"
    if opp >= 7:
        return "med"
    return "low"


def _format_competitor_name(sub: str) -> str:
    """`appstore:Hearth for Contractors` → `Hearth for Contractors (App Store)`."""
    if ":" in sub:
        store, name = sub.split(":", 1)
        store_label = {"appstore": "App Store", "playstore": "Play Store"}.get(store, store)
        return f"{name}  ·  {store_label}"
    return sub


def plan_layout(
    topic: str,
    extra_topics: list[str] | None = None,
    title: str | None = None,
    subtitle: str | None = None,
    tagline: str | None = None,
    max_painpoints: int = 12,
) -> dict[str, Any]:
    """Build a JSON layout plan from corpus data — no rendering.

    Mirrors the schema described in `_doc_design.DESIGN_SYSTEM_PROMPT` so
    an LLM (or a human) can review the structure before committing to
    bytes. Pass the returned plan to `render_planned_docx()` to render.
    """
    data = _gather_deck_data(
        topic, extra_topics=extra_topics, max_painpoints=max_painpoints,
    )
    insights = data.get("insights") or {}

    sources_label = ", ".join(s["source_type"] for s in data["sources"][:6])
    cover = {
        "tagline": (tagline or "Research brief").strip(),
        "title": title or topic,
        "subtitle": subtitle or (
            f"{datetime.now(timezone.utc).strftime('%Y-%m-%d')}  ·  {sources_label}"
        ),
        "kpis": [
            (f"{data['total_posts']:,}", "Corpus posts"),
            (f"{len(data['sources'])}", "Sources"),
            (f"{len(data['painpoints'])}", "Painpoints"),
            (f"{len(data['competitors'])}", "Competitors tracked"),
        ],
    }

    sections: list[dict[str, Any]] = []

    # Executive summary — only render if we have synthesized text
    gt = (insights.get("governing_thought") or "").strip()
    summary = (insights.get("executive_summary") or "").strip()
    if gt or summary:
        sections.append({
            "number": "01",
            "eyebrow": "The thesis",
            "label": "Executive summary",
            "kind": "executive_summary",
            "governing_thought": gt,
            "body": summary,
            "key_arguments": insights.get("key_arguments") or [],
        })

    # Corpus snapshot
    sections.append({
        "number": f"{len(sections)+1:02d}",
        "eyebrow": "What we read",
        "label": "Corpus snapshot",
        "kind": "corpus_table",
        "rows": [(s["source_type"], f"{s['n']:,}") for s in data["sources"][:10]],
        "topics": [topic] + list(extra_topics or []),
    })

    # Painpoint cards
    pps = []
    for i, pp in enumerate(data["painpoints"], 1):
        pps.append({
            "number": i,
            "label": pp["label"],
            "severity": _severity_for_score(pp["opportunity_score"]),
            "frequency": pp["frequency"],
            "opportunity": pp["opportunity_score"],
            "evidence": pp.get("evidence", "").strip(),
            "citations": pp.get("citations") or [],
        })
    sections.append({
        "number": f"{len(sections)+1:02d}",
        "eyebrow": "Where it hurts",
        "label": "Painpoints — evidence-grounded",
        "kind": "painpoint_cards",
        "items": pps,
    })

    # Competitor matrix — sorted by 1★ rate (highest hostility first)
    comps_sorted = sorted(
        data["competitors"],
        key=lambda c: (c["one_star"] / c["reviews"]) if c["reviews"] else 0,
        reverse=True,
    )
    sections.append({
        "number": f"{len(sections)+1:02d}",
        "eyebrow": "Who's already there",
        "label": "Competitor teardown",
        "kind": "competitor_matrix",
        "items": [
            {
                "name": _format_competitor_name(c["sub"]),
                "reviews": c["reviews"],
                "avg_stars": c["avg_stars"],
                "one_star": c["one_star"],
                "five_star": c["five_star"],
                "one_star_rate": (c["one_star"] / c["reviews"]) * 100 if c["reviews"] else 0,
            }
            for c in comps_sorted
        ],
    })

    # Quote wall
    sections.append({
        "number": f"{len(sections)+1:02d}",
        "eyebrow": "Voice of the customer",
        "label": "Top engagement quotes",
        "kind": "quote_wall",
        "items": [
            {
                "quote": (q.get("body") or q.get("title") or "")[:380],
                "title": q.get("title"),
                "attribution": (
                    f"r/{q['sub']}  ·  {q['score']}▲ / {q['num_comments']}\U0001f4ac  ·  id={q['id']}"
                ),
                "permalink": q.get("permalink"),
            }
            for q in data["top_quotes"][:6]
        ],
    })

    # Feature roadmap if we have findings
    findings = sorted(
        insights.get("findings") or [],
        key=lambda f: f.get("opportunity_score", 0),
        reverse=True,
    )[:8]
    if findings:
        sections.append({
            "number": f"{len(sections)+1:02d}",
            "eyebrow": "What we'd build",
            "label": "Top opportunities",
            "kind": "feature_roadmap",
            "items": [
                {
                    "title": f.get("title", ""),
                    "opportunity": f.get("opportunity_score", 0),
                    "narrative": (f.get("narrative") or "").strip(),
                    "best_quote": (f.get("best_quote") or "").strip(),
                }
                for f in findings
            ],
        })

    # Appendix
    sections.append({
        "number": f"{len(sections)+1:02d}",
        "eyebrow": "Reproducible",
        "label": "How to re-pull every citation",
        "kind": "citation_index",
    })

    return {
        "topic": topic,
        "extra_topics": list(extra_topics or []),
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "data_summary": {
            "total_posts": data["total_posts"],
            "sources": len(data["sources"]),
            "painpoints": len(data["painpoints"]),
            "competitors": len(data["competitors"]),
        },
        "cover": cover,
        "sections": sections,
    }


def _render_executive_summary(doc, sec: dict) -> None:
    from . import _doc_design as ds
    if sec.get("governing_thought"):
        ds.add_callout(
            doc, sec["governing_thought"], kind="info", title="The answer",
        )
    if sec.get("body"):
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 1.4
        p.paragraph_format.space_after = Pt(8)
        from ._doc_design import _run, BODY, SIZE_BODY
        _run(p, sec["body"], color=BODY, size=SIZE_BODY)
    args = sec.get("key_arguments") or []
    if args:
        from ._doc_design import _run, INK, SIZE_BODY, SIZE_CAPTION, MUTE, ACCENT
        for i, a in enumerate(args[:3], 1):
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(2)
            _run(p, f"{i:02d}  ", size=Pt(11), color=ACCENT, bold=True)
            _run(p, (a.get("claim") or ""), color=INK, size=SIZE_BODY)


def _render_corpus_table(doc, sec: dict) -> None:
    from . import _doc_design as ds
    rows = sec.get("rows") or []
    if not rows:
        return
    table = doc.add_table(rows=1 + len(rows), cols=2)
    table.autofit = False
    table.rows[0].cells[0].text = "Source"
    table.rows[0].cells[1].text = "Posts"
    for ri, (src, n) in enumerate(rows, 1):
        table.rows[ri].cells[0].text = str(src)
        table.rows[ri].cells[1].text = str(n)
    ds.style_table(table, zebra=True, numeric_cols=[1])
    ds.add_caption(
        doc, "Topics merged: " + " · ".join(sec.get("topics") or []),
    )


def _render_painpoint_cards(doc, sec: dict) -> None:
    from . import _doc_design as ds
    for pp in sec.get("items") or []:
        ds.add_painpoint_card(
            doc,
            number=pp["number"],
            label=pp["label"],
            severity=pp["severity"],
            frequency=pp["frequency"],
            opportunity=pp["opportunity"],
            evidence=pp.get("evidence"),
            citations=pp.get("citations"),
        )


def _render_competitor_matrix(doc, sec: dict) -> None:
    from . import _doc_design as ds
    items = sec.get("items") or []
    if not items:
        return
    cols = ["Competitor", "Reviews", "Avg ★", "1★", "5★", "1★ rate"]
    table = doc.add_table(rows=1 + len(items), cols=len(cols))
    table.autofit = False
    for ci, h in enumerate(cols):
        table.rows[0].cells[ci].text = h
    for ri, c in enumerate(items, 1):
        row = table.rows[ri].cells
        row[0].text = c["name"]
        row[1].text = f"{c['reviews']:,}"
        row[2].text = f"{c['avg_stars']:.2f}"
        row[3].text = f"{c['one_star']:,}"
        row[4].text = f"{c['five_star']:,}"
        row[5].text = f"{c['one_star_rate']:.1f}%"
    ds.style_table(table, zebra=True, numeric_cols=[1, 2, 3, 4, 5])
    # Color-code the 1★ rate column based on hostility threshold.
    for ri, c in enumerate(items, 1):
        cell = table.rows[ri].cells[5]
        rate = c["one_star_rate"]
        if rate >= 30:
            color = ds.PAIN_HI
        elif rate >= 15:
            color = ds.PAIN_MED
        else:
            color = ds.WIN
        for p in cell.paragraphs:
            for r in p.runs:
                r.font.color.rgb = color
                r.bold = True
    ds.add_caption(
        doc, "Sorted by 1★ rate descending — highest-attack-surface competitor at the top.",
    )


def _render_quote_wall(doc, sec: dict) -> None:
    from . import _doc_design as ds
    for q in sec.get("items") or []:
        if q.get("title"):
            from ._doc_design import _para, _run, INK, SIZE_BODY
            p = _para(doc, before=8, after=2)
            _run(p, q["title"], size=SIZE_BODY, color=INK, bold=True)
        ds.add_quote_block(doc, q["quote"], attribution=q["attribution"])


def _render_feature_roadmap(doc, sec: dict) -> None:
    from . import _doc_design as ds
    items = sec.get("items") or []
    if not items:
        return
    cols = ["#", "Opportunity", "Title", "Narrative"]
    table = doc.add_table(rows=1 + len(items), cols=len(cols))
    table.autofit = False
    for ci, h in enumerate(cols):
        table.rows[0].cells[ci].text = h
    for ri, f in enumerate(items, 1):
        row = table.rows[ri].cells
        row[0].text = f"{ri:02d}"
        row[1].text = f"{f.get('opportunity', 0):.1f}/20"
        row[2].text = f.get("title", "")
        row[3].text = (f.get("narrative") or "")[:280]
    ds.style_table(table, zebra=True, numeric_cols=[0, 1])


def _render_citation_index(doc, sec: dict) -> None:
    from . import _doc_design as ds
    from ._doc_design import _run, BODY, SIZE_BODY, FONT_MONO, SIZE_CODE, INK
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(8)
    _run(p, "Every citation in this document is a row in the local SQLite corpus. "
            "Re-read any post with the SQL below — swap <ID> for the post id.",
         color=BODY, size=SIZE_BODY)
    pc = doc.add_paragraph()
    pc.paragraph_format.left_indent = Inches(0.2)
    _run(pc,
         "SELECT id, title, score, sub, source_type, permalink, "
         "substr(selftext,1,1500) AS body FROM posts WHERE id = '<ID>';",
         font=FONT_MONO, size=SIZE_CODE, color=INK)


_RENDERERS = {
    "executive_summary": _render_executive_summary,
    "corpus_table": _render_corpus_table,
    "painpoint_cards": _render_painpoint_cards,
    "competitor_matrix": _render_competitor_matrix,
    "quote_wall": _render_quote_wall,
    "feature_roadmap": _render_feature_roadmap,
    "citation_index": _render_citation_index,
}


def render_planned_docx(plan: dict, out_path: str) -> dict[str, Any]:
    """Render a layout plan to a styled DOCX."""
    if not _DOCX_OK:
        return {
            "ok": False,
            "error": f"python-docx not installed: {_DOCX_ERR}",
            "install_hint": "pip install 'reddit-myind[docs]'",
        }
    from . import _doc_design as ds

    doc = Document()
    # Slightly tighter page margins than Word default — better for printing.
    for section in doc.sections:
        section.top_margin = Inches(0.9)
        section.bottom_margin = Inches(0.9)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    cover = plan.get("cover") or {}
    ds.add_cover_page(
        doc,
        title=cover.get("title", ""),
        subtitle=cover.get("subtitle"),
        tagline=cover.get("tagline"),
        kpis=[tuple(k) for k in (cover.get("kpis") or [])],
    )

    for sec in plan.get("sections") or []:
        ds.add_section_header(
            doc,
            sec.get("label", ""),
            number=sec.get("number"),
            eyebrow=sec.get("eyebrow"),
        )
        renderer = _RENDERERS.get(sec.get("kind"))
        if renderer:
            renderer(doc, sec)
        else:
            from ._doc_design import _run, BODY, SIZE_BODY
            p = doc.add_paragraph()
            _run(p, f"(unsupported section kind: {sec.get('kind')})",
                 color=BODY, italic=True, size=SIZE_BODY)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(out_path)
    return {
        "ok": True,
        "path": out_path,
        "section_count": len(plan.get("sections") or []),
        "output_bytes": Path(out_path).stat().st_size,
    }


def build_docx(
    topic: str,
    out_path: str,
    extra_topics: list[str] | None = None,
    title: str | None = None,
    subtitle: str | None = None,
    tagline: str | None = None,
    max_painpoints: int = 12,
    style: str = "brand",
) -> dict[str, Any]:
    """Render a stakeholder DOCX from corpus data using the brand design system.

    Pipeline: gather → plan → render. The plan is reusable — call
    `plan_layout()` directly if you want to inspect or modify it before
    rendering.

    `style` is a forward-compat hook (currently only "brand"). Future
    values: "minimal" (no cover/no eyebrows), "deck" (one painpoint per
    page).
    """
    if not _DOCX_OK:
        return {
            "ok": False,
            "error": f"python-docx not installed: {_DOCX_ERR}",
            "install_hint": "pip install 'reddit-myind[docs]'",
        }
    plan = plan_layout(
        topic=topic, extra_topics=extra_topics,
        title=title, subtitle=subtitle, tagline=tagline,
        max_painpoints=max_painpoints,
    )
    res = render_planned_docx(plan, out_path)
    if res.get("ok"):
        # Echo the most useful counts so callers see them in the result.
        ds = plan.get("data_summary") or {}
        res.update({
            "engine": "brand-design-system",
            "style": style,
            "painpoint_count": ds.get("painpoints", 0),
            "competitor_count": ds.get("competitors", 0),
            "total_corpus_posts": ds.get("total_posts", 0),
        })
    return res


# ─── PPTX renderer ────────────────────────────────────────────────────────


_TITLE_LAYOUT = 0
_TITLE_AND_CONTENT_LAYOUT = 1
_BLANK_LAYOUT = 6


def _add_title_slide(prs, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_LAYOUT])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_text_slide(prs, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.text = b
        para.font.size = PptxPt(18)


def _add_quote_slide(prs, title: str, quote: str, attribution: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"“{quote}”"
    p.font.size = PptxPt(20)
    p.font.italic = True
    p2 = tf.add_paragraph()
    p2.text = f"— {attribution}"
    p2.font.size = PptxPt(12)
    p2.font.color.rgb = PptxRGBColor(0x55, 0x55, 0x55)


def _add_table_slide(prs, title: str, header: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    title_box = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.6)
    )
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = PptxPt(28)
    tf.paragraphs[0].font.bold = True

    n_rows = len(rows) + 1
    n_cols = len(header)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        PptxInches(0.5), PptxInches(1.1),
        PptxInches(9), PptxInches(min(5.0, 0.4 + 0.32 * n_rows)),
    )
    tbl = tbl_shape.table
    for i, h in enumerate(header):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = PptxPt(12)
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptxPt(10)


def build_pptx(
    topic: str,
    out_path: str,
    extra_topics: list[str] | None = None,
    title: str | None = None,
    subtitle: str | None = None,
    max_painpoints: int = 6,
) -> dict[str, Any]:
    """Render a 12-15 slide pitch deck.

    Slide order:
      1. Cover
      2. TL;DR (governing thought + 3 supporting points)
      3. Corpus snapshot (size + sources)
      4. The market problem (top engagement quote)
      5..N. One slide per top painpoint (label + cited evidence + frequency)
      N+1. Competitor teardown table
      N+2. Voice of the customer (3 best quotes)
      N+3. What we'd build (top features from insights)
      Last. Re-pull instructions

    Returns: {ok, path, slide_count} or {ok: False, error, install_hint}.
    """
    if not _PPTX_OK:
        return {
            "ok": False,
            "error": f"python-pptx not installed: {_PPTX_ERR}",
            "install_hint": "pip install 'reddit-myind[docs]'  # or:  pip install python-pptx",
        }

    data = _gather_deck_data(
        topic, extra_topics=extra_topics, max_painpoints=max_painpoints,
    )
    prs = Presentation()
    insights = data.get("insights") or {}

    # 1. Cover
    _add_title_slide(
        prs,
        title or topic,
        subtitle or f"Research brief · {data['total_posts']} corpus posts · {datetime.now(timezone.utc).strftime('%Y-%m-%d')}",
    )

    # 2. TL;DR
    bullets: list[str] = []
    gt = (insights.get("governing_thought") or "").strip()
    if gt:
        bullets.append(gt)
    for a in (insights.get("key_arguments") or [])[:3]:
        c = (a.get("claim") or "").strip()
        if c:
            bullets.append(c)
    if not bullets:
        bullets = [
            f"{data['total_posts']} corpus posts across {len(data['sources'])} sources.",
            f"Top painpoints derived from {len(data['painpoints'])} extracted graph nodes.",
            f"Competitor signal from {len(data['competitors'])} apps with ≥5 reviews each.",
        ]
    _add_text_slide(prs, "TL;DR", bullets)

    # 3. Corpus snapshot
    rows = [[s["source_type"], str(s["n"])] for s in data["sources"][:10]]
    _add_table_slide(prs, "Corpus snapshot", ["Source", "Posts"], rows)

    # 4. Market problem (top quote)
    if data["top_quotes"]:
        q = data["top_quotes"][0]
        _add_quote_slide(
            prs,
            "The market problem (top engagement)",
            (q.get("body") or q.get("title") or "")[:300],
            f"r/{q['sub']} · {q['score']}▲ / {q['num_comments']}\U0001f4ac · id={q['id']}",
        )

    # 5..N. One slide per painpoint
    for i, pp in enumerate(data["painpoints"], 1):
        attribution_bits = [
            f"Frequency {pp['frequency']}",
            f"Opportunity {pp['opportunity_score']:.1f}/20",
        ]
        if pp.get("citations"):
            c = pp["citations"][0]
            attribution_bits.append(f"e.g. {c['source_type']}/{c['sub']} · id={c['id']}")
        _add_quote_slide(
            prs,
            f"Painpoint {i}: {pp['label']}",
            (pp.get("evidence") or "(no extracted quote)")[:280],
            " · ".join(attribution_bits),
        )

    # N+1. Competitor matrix
    if data["competitors"]:
        comp_rows = []
        for c in data["competitors"]:
            rate = (c["one_star"] / c["reviews"]) * 100 if c["reviews"] else 0
            comp_rows.append([
                c["sub"], str(c["reviews"]), f"{c['avg_stars']:.2f}",
                str(c["one_star"]), f"{rate:.0f}%",
            ])
        _add_table_slide(
            prs, "Competitor teardown",
            ["Competitor", "Reviews", "Avg ★", "1★", "1★ rate"],
            comp_rows,
        )

    # N+2. Voice of the customer
    if len(data["top_quotes"]) > 1:
        for q in data["top_quotes"][1:4]:
            _add_quote_slide(
                prs,
                "Voice of the customer",
                (q.get("body") or q.get("title") or "")[:260],
                f"r/{q['sub']} · {q['score']}▲ · id={q['id']}",
            )

    # N+3. What we'd build (from synthesized findings if available)
    findings = sorted(
        insights.get("findings") or [],
        key=lambda f: f.get("opportunity_score", 0),
        reverse=True,
    )[:6]
    if findings:
        bullets = [f"{f.get('title','')} (opp {f.get('opportunity_score', 0):.1f}/20)" for f in findings]
        _add_text_slide(prs, "What we'd build (top opportunities)", bullets)

    # Last: re-pull
    _add_text_slide(prs, "Every citation is re-pullable", [
        "Citations carry source_type / sub / id.",
        "SELECT * FROM posts WHERE id = '<ID>';",
        "App Store / Play Store IDs include the app name in `sub`.",
        "Reddit IDs map to https://www.reddit.com/r/<sub>/comments/<id>/.",
    ])

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return {
        "ok": True,
        "path": out_path,
        "slide_count": len(prs.slides),
        "painpoint_count": len(data["painpoints"]),
        "competitor_count": len(data["competitors"]),
        "total_corpus_posts": data["total_posts"],
    }


def _ensure_brand_reference_docx() -> str | None:
    """Return path to the cached brand reference docx; generate on first use.

    Falls back silently if python-docx isn't available — pandoc will
    just use its default styles.
    """
    if not _DOCX_OK:
        return None
    try:
        from ..core.config import _resolve_data_dir
        base = _resolve_data_dir()
    except Exception:
        base = Path.home() / ".gapmap"
        base.mkdir(parents=True, exist_ok=True)
    cache = Path(base) / "brand-reference.docx"
    if cache.exists():
        return str(cache)
    try:
        from . import _doc_design as ds
        return ds.make_brand_reference_docx(cache)
    except Exception:
        return None


def build_docx_from_markdown(
    md_path: str,
    out_path: str,
    reference_docx: str | None = None,
) -> dict[str, Any]:
    """Convert a markdown research brief to DOCX with full fidelity.

    Use this when you've already authored the rich research doc as
    markdown (cited evidence quotes, competitor tables, code blocks,
    blockquotes, headings) and want the same content as a Word file
    that opens cleanly in Word / Pages / Google Docs.

    Strategy:
      - **Primary:** pandoc via `pypandoc`. Pandoc is the gold standard
        for markdown → docx — preserves headings, GitHub-flavored tables,
        blockquotes, fenced code, bold/italic, links, lists, footnotes.
      - **Fallback:** if pandoc isn't available, fall back to a
        line-by-line python-docx renderer that handles headings, lists,
        blockquotes, tables, and bold/italic inline runs. Lower fidelity
        on edge cases but ships zero-extra-deps.

    `reference_docx` (optional): a Word doc whose styles/fonts pandoc
    will copy (cover page, headers, fonts, accent colors). Skip for the
    default Word styles.

    Returns: {ok, path, engine, source_chars, output_bytes} or {ok: False, error}.
    """
    src = Path(md_path)
    if not src.exists():
        return {"ok": False, "error": f"Markdown source not found: {md_path}"}

    md_text = src.read_text(encoding="utf-8")
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)

    # --- Primary path: pandoc ---
    # Auto-generate the brand reference docx so markdown→docx output
    # carries the same fonts/colors/spacing as the data-driven path.
    if not reference_docx:
        reference_docx = _ensure_brand_reference_docx()

    if _PANDOC_OK:
        try:
            extra_args = [
                "--standalone",
                "--toc",
                "--toc-depth=2",
                # GFM tables + auto-link bare URLs + strikethrough — matches
                # the markdown the deep-dive doc was authored in.
                "--from=gfm+pipe_tables+task_lists+autolink_bare_uris+strikeout",
            ]
            if reference_docx and Path(reference_docx).exists():
                extra_args += ["--reference-doc", reference_docx]
            pypandoc.convert_text(
                md_text,
                to="docx",
                format="gfm",
                outputfile=out_path,
                extra_args=extra_args,
            )
            return {
                "ok": True,
                "path": out_path,
                "engine": "pandoc",
                "source_chars": len(md_text),
                "output_bytes": Path(out_path).stat().st_size,
            }
        except Exception as e:
            # Fall through to the python-docx renderer with a note.
            pandoc_err = str(e)
    else:
        pandoc_err = _PANDOC_ERR if not _PANDOC_OK else None

    # --- Fallback: python-docx renderer ---
    if not _DOCX_OK:
        return {
            "ok": False,
            "error": f"pandoc unavailable ({pandoc_err}) and python-docx not installed: {_DOCX_ERR}",
            "install_hint": "pip install 'reddit-myind[docs]'",
        }

    doc = Document()
    _render_markdown_to_docx(doc, md_text)
    doc.save(out_path)
    return {
        "ok": True,
        "path": out_path,
        "engine": "python-docx-fallback",
        "pandoc_unavailable_reason": pandoc_err,
        "source_chars": len(md_text),
        "output_bytes": Path(out_path).stat().st_size,
    }


# ─── Lightweight markdown → python-docx fallback renderer ────────────────


def _render_inline(paragraph, text: str) -> None:
    """Apply **bold**, *italic*, _italic_, `code` inline runs to a paragraph."""
    import re
    # Pattern order matters: triple before double before single.
    pattern = re.compile(
        r"(\*\*\*[^*]+\*\*\*"     # ***bold-italic***
        r"|\*\*[^*]+\*\*"         # **bold**
        r"|__[^_]+__"             # __bold__
        r"|\*[^*]+\*"             # *italic*
        r"|_[^_]+_"               # _italic_
        r"|`[^`]+`"               # `code`
        r")"
    )
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            paragraph.add_run(text[pos:m.start()])
        token = m.group(1)
        run = None
        if token.startswith("***") and token.endswith("***"):
            run = paragraph.add_run(token[3:-3])
            run.bold = True
            run.italic = True
        elif (token.startswith("**") and token.endswith("**")) or (
            token.startswith("__") and token.endswith("__")
        ):
            run = paragraph.add_run(token[2:-2])
            run.bold = True
        elif (token.startswith("*") and token.endswith("*")) or (
            token.startswith("_") and token.endswith("_")
        ):
            run = paragraph.add_run(token[1:-1])
            run.italic = True
        elif token.startswith("`") and token.endswith("`"):
            run = paragraph.add_run(token[1:-1])
            run.font.name = "Menlo"
            run.font.size = Pt(10)
        pos = m.end()
    if pos < len(text):
        paragraph.add_run(text[pos:])


def _render_markdown_to_docx(doc, md_text: str) -> None:
    """Streaming markdown → docx for the doc structure used by this repo:

      - ATX headings (#, ##, …)
      - Fenced code blocks (```)
      - Pipe tables (with optional alignment row)
      - Blockquotes (>)
      - Bullet/numbered lists
      - Bold/italic/code inline runs
      - Horizontal rules

    Designed for *our* markdown — not a general parser. Edge cases like
    nested lists or footnotes degrade to plain paragraphs.
    """
    import re

    lines = md_text.splitlines()
    i = 0
    in_code = False
    code_buf: list[str] = []

    while i < len(lines):
        line = lines[i]

        # Fenced code block
        if line.lstrip().startswith("```"):
            if in_code:
                p = doc.add_paragraph()
                run = p.add_run("\n".join(code_buf))
                run.font.name = "Menlo"
                run.font.size = Pt(9)
                code_buf = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_buf.append(line)
            i += 1
            continue

        # Horizontal rule
        if re.match(r"^\s*-{3,}\s*$", line) or re.match(r"^\s*\*{3,}\s*$", line):
            doc.add_paragraph("―" * 30).alignment = WD_ALIGN_PARAGRAPH.CENTER
            i += 1
            continue

        # ATX headings
        m = re.match(r"^(#{1,6})\s+(.+?)\s*#*\s*$", line)
        if m:
            level = min(len(m.group(1)), 4)
            doc.add_heading(m.group(2), level=level)
            i += 1
            continue

        # Pipe table — first row is header, second row is alignment, then body
        if line.lstrip().startswith("|") and i + 1 < len(lines) and re.match(
            r"^\s*\|?\s*:?-{2,}", lines[i + 1]
        ):
            tbl_lines = [line]
            j = i + 1
            while j < len(lines) and lines[j].lstrip().startswith("|"):
                tbl_lines.append(lines[j])
                j += 1

            def split_row(s: str) -> list[str]:
                s = s.strip()
                if s.startswith("|"):
                    s = s[1:]
                if s.endswith("|"):
                    s = s[:-1]
                return [c.strip() for c in s.split("|")]

            header = split_row(tbl_lines[0])
            body = [split_row(r) for r in tbl_lines[2:]]
            n_cols = len(header)
            if n_cols >= 1:
                table = doc.add_table(rows=1 + len(body), cols=n_cols)
                table.style = "Light Grid Accent 1"
                for ci, h in enumerate(header):
                    cell = table.rows[0].cells[ci]
                    cell.text = ""
                    p = cell.paragraphs[0]
                    run = p.add_run(h)
                    run.bold = True
                for ri, row in enumerate(body, 1):
                    for ci in range(n_cols):
                        val = row[ci] if ci < len(row) else ""
                        cell = table.rows[ri].cells[ci]
                        cell.text = ""
                        _render_inline(cell.paragraphs[0], val)
                doc.add_paragraph()
            i = j
            continue

        # Blockquote
        if line.lstrip().startswith(">"):
            buf = []
            while i < len(lines) and lines[i].lstrip().startswith(">"):
                buf.append(re.sub(r"^\s*>\s?", "", lines[i]))
                i += 1
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.4)
            _render_inline(p, " ".join(buf))
            for run in p.runs:
                run.italic = True
            continue

        # Numbered list
        m = re.match(r"^\s*(\d+)\.\s+(.+)$", line)
        if m:
            p = doc.add_paragraph(style="List Number")
            _render_inline(p, m.group(2))
            i += 1
            continue

        # Bullet list
        m = re.match(r"^\s*[-*+]\s+(.+)$", line)
        if m:
            p = doc.add_paragraph(style="List Bullet")
            _render_inline(p, m.group(1))
            i += 1
            continue

        # Blank line — paragraph separator (already handled by add_paragraph).
        if not line.strip():
            i += 1
            continue

        # Default: paragraph with inline formatting
        p = doc.add_paragraph()
        _render_inline(p, line)
        i += 1

    # Flush any unterminated code block
    if in_code and code_buf:
        p = doc.add_paragraph()
        run = p.add_run("\n".join(code_buf))
        run.font.name = "Menlo"
        run.font.size = Pt(9)


def get_design_system_prompt() -> str:
    """Strict design-rules prompt for an LLM that's planning a layout."""
    from ._doc_design import DESIGN_SYSTEM_PROMPT
    return DESIGN_SYSTEM_PROMPT


# ─── PDF builder (markdown → polished PDF via pandoc + xelatex) ──────────
# Mirrors the docs/demo_pdf/pdf_build pipeline that produced the
# lending-marketplace deep-dive. The custom LaTeX header (header.tex) and
# the citation-table widening Lua filter (widen-quote.lua) are bundled
# inside the package under research/_doc_assets/pdf/, so wheel users get
# the same brand without copying files around.
#
# Requires: pypandoc + xelatex on PATH. We auto-detect xelatex; if it's
# missing we return a structured error explaining how to install BasicTeX
# or MacTeX. All optional — the function never crashes the MCP server.


def _find_xelatex() -> str | None:
    """Best-effort discovery of xelatex on PATH (or in the standard
    MacTeX/TeX-Live install locations).
    """
    import shutil
    found = shutil.which("xelatex")
    if found:
        return found
    for candidate in (
        "/Library/TeX/texbin/xelatex",
        "/usr/local/texlive/2024/bin/universal-darwin/xelatex",
        "/usr/local/texlive/2023/bin/universal-darwin/xelatex",
        "/usr/local/bin/xelatex",
        "/opt/homebrew/bin/xelatex",
    ):
        if Path(candidate).exists():
            return candidate
    return None


def _pdf_assets_dir() -> Path:
    """Return the bundled-assets directory shipped with this package."""
    return Path(__file__).parent / "_doc_assets" / "pdf"


def build_pdf_from_markdown(
    md_path: str,
    out_path: str,
    title: str | None = None,
    subtitle: str | None = None,
    brand_link: str = "gapmap.myind.ai",
    brand_link_url: str = "https://gapmap.myind.ai",
    extra_pandoc_args: list[str] | None = None,
) -> dict[str, Any]:
    """Convert a markdown research brief to a brand-styled PDF.

    Pipeline (matches `docs/demo_pdf/pdf_build/`):
      pandoc <md> --pdf-engine=xelatex
                  --include-in-header=<bundled header.tex>
                  --lua-filter=<bundled widen-quote.lua>
                  --variable=mainfont:"DejaVu Sans" \\
                  --variable=sansfont:"Poppins" \\
                  --variable=monofont:"DejaVu Sans Mono" \\
                  --variable=geometry:a4paper,top=2.4cm,bottom=2.4cm,left=2.0cm,right=2.0cm \\
                  --variable=fontsize:10pt \\
                  --toc --toc-depth=2

    `title` and `subtitle` populate the `\\brandtitle{}` and
    `\\brandsubtitle{}` macros used by fancyhdr. Defaults read from the
    markdown's first H1.

    Returns: {ok, path, engine: 'xelatex', source_chars, output_bytes}
    or a structured error with `install_hint` if pandoc / xelatex is
    missing.
    """
    src = Path(md_path)
    if not src.exists():
        return {"ok": False, "error": f"Markdown source not found: {md_path}"}

    if not _PANDOC_OK:
        return {
            "ok": False,
            "error": f"pypandoc not installed: {_PANDOC_ERR}",
            "install_hint": "pip install 'reddit-myind[docs]'  # bundles pandoc binary",
        }

    xelatex = _find_xelatex()
    if not xelatex:
        return {
            "ok": False,
            "error": "xelatex not found on PATH",
            "install_hint": (
                "Install BasicTeX (~80 MB) or MacTeX (~4 GB):\n"
                "  brew install --cask basictex\n"
                "  sudo tlmgr update --self && sudo tlmgr install xetex titlesec "
                "titling fancyhdr enumitem microtype fvextra quoting seqsplit"
            ),
        }

    assets = _pdf_assets_dir()
    header_tex = assets / "header.tex"
    lua_filter = assets / "widen-quote.lua"
    if not header_tex.exists() or not lua_filter.exists():
        return {
            "ok": False,
            "error": (
                f"Bundled PDF assets missing under {assets}. "
                "Reinstall the package or copy header.tex + widen-quote.lua "
                "from docs/demo_pdf/pdf_build/."
            ),
        }

    md_text = src.read_text(encoding="utf-8")

    # If the caller didn't pass a title, lift the first H1.
    auto_title = title
    if not auto_title:
        for line in md_text.splitlines():
            if line.startswith("# "):
                auto_title = line[2:].strip()
                break

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)

    extra_args = [
        "--standalone",
        "--toc",
        "--toc-depth=2",
        "--from=gfm+pipe_tables+task_lists+autolink_bare_uris+strikeout",
        "--pdf-engine", xelatex,
        f"--include-in-header={header_tex}",
        f"--lua-filter={lua_filter}",
        # Geometry + fonts — same as the demo_pdf build.
        "--variable=geometry:a4paper,top=2.4cm,bottom=2.4cm,left=2.0cm,right=2.0cm",
        "--variable=fontsize:10pt",
        "--variable=mainfont:DejaVu Sans",
        "--variable=sansfont:Poppins",
        "--variable=monofont:DejaVu Sans Mono",
        # Brand chrome strings the header.tex picks up via providecommand.
        f"--variable=brandtitle:{auto_title or 'Research brief'}",
        f"--variable=brandsubtitle:{subtitle or 'Gap Map'}",
        f"--variable=brandlink:{brand_link}",
        f"--variable=brandlinkurl:{brand_link_url}",
    ]
    if extra_pandoc_args:
        extra_args.extend(extra_pandoc_args)

    try:
        pypandoc.convert_text(
            md_text,
            to="pdf",
            format="gfm",
            outputfile=out_path,
            extra_args=extra_args,
        )
    except Exception as e:
        return {
            "ok": False,
            "error": f"pandoc/xelatex failed: {e}",
            "install_hint": (
                "Most LaTeX errors are missing packages. Run:\n"
                "  sudo tlmgr install titlesec titling fancyhdr enumitem "
                "microtype fvextra quoting seqsplit framed"
            ),
        }

    out = Path(out_path)
    return {
        "ok": True,
        "path": str(out),
        "engine": "xelatex",
        "source_chars": len(md_text),
        "output_bytes": out.stat().st_size if out.exists() else 0,
        "header_tex": str(header_tex),
        "lua_filter": str(lua_filter),
    }


__all__ = [
    "build_docx", "build_pptx", "build_docx_from_markdown",
    "build_pdf_from_markdown",
    "plan_layout", "render_planned_docx", "get_design_system_prompt",
]
